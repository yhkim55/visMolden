from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_centered_text(slide, left, top, width, height, text, font_size=10):
    tx_box = slide.shapes.add_textbox(left, top, width, height)

    p = tx_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(0, 0, 0)


def collate_png_ppt(
    png_map,
    occ,
    orb_indices,
    ppt_fn,
):
    """
    Make PPT slides with a 3x5 grid of orbital images.

    Parameters
    ----------
    png_map : dict[int, str]
        Orbital index -> PNG path
    occ : ndarray
        Occupation array
    orb_indices : list[int]
        Orbitals to visualize
    ppt_fn : str
        Output PPT filename
    """
    prs = Presentation()

    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    layout = prs.slide_layouts[6]  # blank slide

    ncol = 5
    nrow = 3
    per_slide = ncol * nrow

    img_w = Inches(1.9)
    txt_h = Inches(0.25)

    margin_x = Inches(0.25)
    margin_y = Inches(0.25)

    gap_x = Inches(0.15)
    gap_y = Inches(0.35)

    for start in range(0, len(orb_indices), per_slide):

        slide = prs.slides.add_slide(layout)

        subset = orb_indices[start:start + per_slide]

        for k, idx in enumerate(subset):

            row = k // ncol
            col = k % ncol

            left = margin_x + col * (img_w + gap_x)

            top = margin_y + row * (img_w + txt_h + gap_y)

            slide.shapes.add_picture(
                png_map[idx],
                left,
                top,
                width=img_w,
            )

            add_centered_text(
                slide,
                left,
                top + img_w,
                img_w,
                txt_h,
                f"idx={idx}  occ={occ[idx]:.3f}",
                font_size=10,
            )

    prs.save(ppt_fn)